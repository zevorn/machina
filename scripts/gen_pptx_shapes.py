"""
Generate os2atc-2026.pptx with native PowerPoint shapes
instead of ASCII art code blocks.
"""

import os

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# -- constants --
SLIDE_W = 12192000
SLIDE_H = 6858000
ACCENT = RGBColor(0xE2, 0x6C, 0x0D)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
BG_ACCENT = RGBColor(0xFD, 0xF0, 0xE5)
BG_CODE = RGBColor(0xF8, 0xF8, 0xF8)
FONT_CN = "微软雅黑"
FONT_CODE = "Consolas"
FOOTER_TEXT = "OS2ATC 2026 · AI 辅助编程分论坛"
TOTAL_SLIDES = 18
_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "assets")
IMG_BG = os.environ.get("PPTX_BG",
                        os.path.join(_ASSETS, "bg.png"))
IMG_LOGO = os.environ.get("PPTX_LOGO",
                          os.path.join(_ASSETS, "logo.png"))
IMG_LOGO2 = os.environ.get("PPTX_LOGO2",
                           os.path.join(_ASSETS, "logo2.png"))

E = Emu  # shorthand

_img_warned = set()


def safe_add_picture(sl, path, left, top, width, height=None):
    """Add picture if file exists, skip with warning otherwise."""
    if not os.path.exists(path):
        if path not in _img_warned:
            print(f"  warn: {path} not found, skipping")
            _img_warned.add(path)
        return
    if height is not None:
        sl.shapes.add_picture(path, left, top, width, height)
    else:
        sl.shapes.add_picture(path, left, top, width)


# ============================================================
# Shape helpers
# ============================================================

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_run(p, text, font_name=FONT_CN, size=Pt(14),
            bold=False, color=DARK, italic=False):
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return r


def add_textbox(sl, left, top, width, height):
    return sl.shapes.add_textbox(E(left), E(top),
                                 E(width), E(height))


def add_accent_bar(sl):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            0, E(324000), E(108000), E(504000))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def add_title(sl, text):
    add_accent_bar(sl)
    t = add_textbox(sl, 252000, 270000, 9000000, 576000)
    set_run(t.text_frame.paragraphs[0], text,
            size=Pt(26), color=ACCENT)


def add_footer(sl, n, total):
    t = add_textbox(sl, 540000, 6480000, 6480000, 252000)
    set_run(t.text_frame.paragraphs[0], FOOTER_TEXT,
            size=Pt(11), color=GRAY)
    t2 = add_textbox(sl, 10800000, 6480000, 1080000, 252000)
    t2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_run(t2.text_frame.paragraphs[0], f"{n}/{total}",
            size=Pt(11), color=GRAY)


def add_logo(sl):
    safe_add_picture(sl, IMG_LOGO2,
                     E(10260000), E(108000),
                     E(1368000), E(570610))


def add_highlight_box(sl, text, left=540000, top=5400000,
                      width=10800000, height=540000):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            E(left), E(top), E(width), E(height))
    s.fill.solid()
    s.fill.fore_color.rgb = BG_ACCENT
    s.line.color.rgb = ACCENT
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = E(180000)
    tf.margin_top = E(72000)
    set_run(tf.paragraphs[0], text,
            size=Pt(14), bold=True, color=ACCENT)


def add_quote(sl, text, left=540000, top=5400000,
              width=10800000, height=500000):
    b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            E(left), E(top), E(54000), E(height))
    b.fill.solid()
    b.fill.fore_color.rgb = ACCENT
    b.line.fill.background()
    t = add_textbox(sl, left + 108000, top,
                    width - 108000, height)
    t.text_frame.word_wrap = True
    set_run(t.text_frame.paragraphs[0], text,
            size=Pt(13), italic=True, color=GRAY)


def add_code_block(sl, lines, left=540000, top=1008000,
                   width=10800000, height=5200000,
                   font_size=Pt(12)):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            E(left), E(top), E(width), E(height))
    s.fill.solid()
    s.fill.fore_color.rgb = BG_CODE
    s.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = E(180000)
    tf.margin_top = E(144000)
    tf.margin_right = E(180000)
    tf.margin_bottom = E(144000)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        set_run(p, line, font_name=FONT_CODE,
                size=font_size, color=DARK)


def add_table(sl, headers, rows, left=540000, top=1008000,
              width=10800000, row_height=360000):
    tr = 1 + len(rows)
    s = sl.shapes.add_table(tr, len(headers),
                            E(left), E(top),
                            E(width), E(tr * row_height))
    t = s.table
    cw = width // len(headers)
    for i in range(len(headers)):
        t.columns[i].width = E(cw)
    for ci, h in enumerate(headers):
        c = t.cell(0, ci)
        c.text = ""
        set_run(c.text_frame.paragraphs[0], h,
                size=Pt(14), bold=True, color=WHITE)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
    for ri, r in enumerate(rows):
        for ci, v in enumerate(r):
            c = t.cell(ri + 1, ci)
            c.text = ""
            set_run(c.text_frame.paragraphs[0], v,
                    size=Pt(13), color=DARK)
            c.fill.solid()
            c.fill.fore_color.rgb = (
                BG_LIGHT if ri % 2 == 0 else WHITE)


# -- shape diagram primitives --

def draw_box(sl, left, top, width, height, lines,
             fill=BG_LIGHT, border=RGBColor(0xA5, 0xA5, 0xA5),
             font_size=Pt(12), bold=False, text_color=DARK,
             align=PP_ALIGN.CENTER):
    """Draw a rounded rectangle with centered multi-line text."""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            E(left), E(top), E(width), E(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = E(72000)
    tf.margin_right = E(72000)
    tf.margin_top = E(36000)
    tf.margin_bottom = E(36000)
    # vertical center
    try:
        tf.paragraphs[0].alignment = align
        txBody = tf._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('anchor', 'ctr')
    except Exception:
        pass
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(1)
        set_run(p, line, size=font_size, bold=bold,
                color=text_color)
    return s


def draw_arrow_right(sl, x1, y, x2, color=ACCENT, width=Pt(2)):
    """Horizontal arrow from (x1, y) to (x2, y)."""
    connector = sl.shapes.add_connector(
        1, E(x1), E(y), E(x2), E(y))  # type 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    line = connector.line
    line_elem = line._ln
    tailEnd = line_elem.makeelement(qn('a:tailEnd'), {})
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')
    line_elem.append(tailEnd)
    return connector


def draw_arrow_down(sl, x, y1, y2, color=ACCENT, width=Pt(2)):
    """Vertical arrow from (x, y1) to (x, y2)."""
    connector = sl.shapes.add_connector(
        1, E(x), E(y1), E(x), E(y2))
    connector.line.color.rgb = color
    connector.line.width = width
    line_elem = connector.line._ln
    tailEnd = line_elem.makeelement(qn('a:tailEnd'), {})
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')
    line_elem.append(tailEnd)
    return connector


def draw_arrow_left(sl, x1, y, x2, color=ACCENT, width=Pt(2)):
    """Horizontal arrow from (x1, y) to (x2, y), pointing left."""
    connector = sl.shapes.add_connector(
        1, E(x1), E(y), E(x2), E(y))
    connector.line.color.rgb = color
    connector.line.width = width
    line_elem = connector.line._ln
    # head on the x2 side (tail)
    tailEnd = line_elem.makeelement(qn('a:tailEnd'), {})
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')
    line_elem.append(tailEnd)
    return connector


def draw_line(sl, x1, y1, x2, y2, color=ACCENT, width=Pt(1.5)):
    """Plain line without arrowhead."""
    c = sl.shapes.add_connector(1, E(x1), E(y1), E(x2), E(y2))
    c.line.color.rgb = color
    c.line.width = width
    return c


def draw_label(sl, left, top, width, height, text,
               font_size=Pt(11), color=GRAY, bold=False,
               align=PP_ALIGN.CENTER):
    """Small text label."""
    t = add_textbox(sl, left, top, width, height)
    t.text_frame.paragraphs[0].alignment = align
    set_run(t.text_frame.paragraphs[0], text,
            size=font_size, color=color, bold=bold)
    return t


# ============================================================
# Slide builders
# ============================================================

def slide_pipeline(sl):
    """Slide 3: Translation Pipeline with shapes."""
    # Layout: 4 boxes top row, 3 boxes bottom row + TB Cache
    BW = 2200000   # box width
    BH = 750000    # box height
    GAP = 300000   # horizontal gap
    TOP = 1100000
    LEFT = 540000

    def bx(col):
        return LEFT + col * (BW + GAP)

    # Top row
    boxes_top = [
        (["Guest ELF", "(RISC-V)"], BG_LIGHT),
        (["Frontend", "Decode"], BG_LIGHT),
        (["TCG IR", "(Intermediate)"], BG_LIGHT),
        (["Optimizer", "ConstFold + CopyProp"], BG_ACCENT),
    ]
    for i, (lines, fill) in enumerate(boxes_top):
        x = bx(i)
        border = ACCENT if fill == BG_ACCENT else \
            RGBColor(0xA5, 0xA5, 0xA5)
        draw_box(sl, x, TOP, BW, BH, lines,
                 fill=fill, border=border, font_size=Pt(12))
        if i < 3:
            draw_arrow_right(sl, x + BW, TOP + BH // 2,
                             bx(i + 1))

    # Bottom row
    BOT = TOP + BH + 600000
    boxes_bot = [
        (["Result"], BG_LIGHT),
        (["JIT Exec"], BG_LIGHT),
        (["Backend", "x86-64 Codegen"], BG_LIGHT),
    ]
    for i, (lines, fill) in enumerate(boxes_bot):
        x = bx(i)
        draw_box(sl, x, BOT, BW, BH, lines,
                 fill=fill, font_size=Pt(12))
        if i < 2:
            draw_arrow_left(sl, bx(i + 1), BOT + BH // 2,
                            x + BW)

    # Arrow: Optimizer down to Backend (center-top)
    opt_x = bx(3) + BW // 2
    back_cx = bx(2) + BW // 2
    mid_y = BOT - 150000
    draw_line(sl, opt_x, TOP + BH,
              opt_x, mid_y)
    draw_line(sl, opt_x, mid_y,
              back_cx, mid_y)
    draw_arrow_down(sl, back_cx, mid_y, BOT)

    # TB Cache box
    TB_W = 2000000
    TB_H = 500000
    tb_x = bx(1) + (BW - TB_W) // 2
    tb_y = BOT + BH + 400000
    draw_box(sl, tb_x, tb_y, TB_W, TB_H,
             ["TB Cache"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(12), bold=True)
    # Arrow from JIT up to TB Cache (or label)
    jit_cx = bx(1) + BW // 2
    draw_arrow_down(sl, jit_cx, BOT + BH, tb_y)

    draw_label(sl, tb_x + TB_W + 100000, tb_y,
               3000000, TB_H,
               "translate once,\nexecute many",
               font_size=Pt(12), color=ACCENT, bold=True,
               align=PP_ALIGN.LEFT)


def slide_roles(sl):
    """Slide 4: Role Division with 3 columns."""
    COL_W = 3300000
    GAP = 200000
    LEFT = 540000
    TOP = 1050000
    HDR_H = 700000
    BODY_H = 3200000

    def cx(i):
        return LEFT + i * (COL_W + GAP)

    titles = [
        ("Human", "Product Manager"),
        ("Claude", "Architect + Lead Dev"),
        ("Codex", "Review + Completion"),
    ]
    fills = [BG_LIGHT, BG_ACCENT, BG_LIGHT]
    borders = [RGBColor(0xA5, 0xA5, 0xA5), ACCENT,
               RGBColor(0xA5, 0xA5, 0xA5)]
    for i, ((t1, t2), fill, brd) in enumerate(
            zip(titles, fills, borders)):
        draw_box(sl, cx(i), TOP, COL_W, HDR_H,
                 [t1, t2], fill=fill, border=brd,
                 font_size=Pt(14), bold=True)

    bodies = [
        ["· Define WHAT", "· Define WHY",
         "· Verify design doc", "· Verify test result",
         "· Perf direction"],
        ["· Plan Mode", "  Decide & Plan",
         "· Core logic impl", "· Write design doc"],
        ["· Code Review", "· Batch test",
         "· Cross-verify", "· Insn completion"],
    ]
    BODY_TOP = TOP + HDR_H + 50000
    for i, lines in enumerate(bodies):
        draw_box(sl, cx(i), BODY_TOP, COL_W, BODY_H,
                 lines, fill=WHITE, border=borders[i],
                 font_size=Pt(12),
                 align=PP_ALIGN.LEFT)


def slide_decisions(sl):
    """Slide 5: Plan Mode decision flow."""
    LEFT = 540000
    # Info box at top
    draw_box(sl, LEFT, 1050000, 10800000, 700000,
             ["Human experience guides Agent's dev path",
              "NOT top-down design → AI fills."
              "  Follow human's path: backend first"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(13), align=PP_ALIGN.LEFT)

    # 4 step boxes
    BW = 2300000
    BH = 800000
    GAP = 300000
    TOP2 = 2100000

    steps = [
        ("1. Analyze", "x86 insn format"),
        ("2. Plan", "Prioritize steps"),
        ("3. Discuss", "Adjust & refine"),
        ("4. Execute", "Implement"),
    ]
    for i, (t1, t2) in enumerate(steps):
        x = LEFT + i * (BW + GAP)
        fill = BG_ACCENT if i == 3 else BG_LIGHT
        brd = ACCENT if i == 3 else RGBColor(0xA5, 0xA5, 0xA5)
        draw_box(sl, x, TOP2, BW, BH, [t1, t2],
                 fill=fill, border=brd, font_size=Pt(13))
        if i < 3:
            draw_arrow_right(sl, x + BW, TOP2 + BH // 2,
                             x + BW + GAP)

    # "Agreed" label between step 3 and 4
    draw_label(sl, LEFT + 2 * (BW + GAP) + BW,
               TOP2 - 200000, GAP, 200000,
               "Agreed", font_size=Pt(10), color=ACCENT)


def slide_doc_driven(sl):
    """Slide 6: Design doc iteration cycle."""
    BW = 1800000
    BH = 650000
    GAP = 250000
    LEFT = 540000
    TOP = 1100000

    # 5 boxes in a cycle:
    # 1.Define -> 2.Write doc -> 3.Generate code
    #    ^                          |
    #    +---- 5.Update doc <-- 4.Diff test
    boxes = [
        (LEFT, TOP, "1. Define\nRequire"),
        (LEFT + BW + GAP, TOP, "2. Agent\nWrite design doc"),
        (LEFT + 2 * (BW + GAP), TOP,
         "3. Agent\nGenerate code"),
        (LEFT + 2 * (BW + GAP), TOP + BH + 500000,
         "4. Diff test\nVerify"),
        (LEFT + BW + GAP, TOP + BH + 500000,
         "5. Update\nDesign doc"),
    ]
    for i, (x, y, text) in enumerate(boxes):
        fill = BG_ACCENT if i in (1, 4) else BG_LIGHT
        brd = ACCENT if i in (1, 4) else \
            RGBColor(0xA5, 0xA5, 0xA5)
        draw_box(sl, x, y, BW, BH, text.split("\n"),
                 fill=fill, border=brd, font_size=Pt(12))

    # Arrows: 1->2
    draw_arrow_right(sl, LEFT + BW,
                     TOP + BH // 2,
                     LEFT + BW + GAP)
    # 2->3
    draw_arrow_right(sl, LEFT + BW + GAP + BW,
                     TOP + BH // 2,
                     LEFT + 2 * (BW + GAP))
    # 3->4 (down)
    x3 = LEFT + 2 * (BW + GAP) + BW // 2
    draw_arrow_down(sl, x3, TOP + BH,
                    TOP + BH + 500000)
    # 4->5 (left)
    x5r = LEFT + BW + GAP + BW
    x4l = LEFT + 2 * (BW + GAP)
    y45 = TOP + BH + 500000 + BH // 2
    draw_arrow_left(sl, x4l, y45, x5r)
    # 5->1 (left + up, routed outside box 1)
    x5l = LEFT + BW + GAP
    y5 = TOP + BH + 500000 + BH // 2
    y1 = TOP + BH // 2
    feedback_x = LEFT - 200000
    draw_line(sl, x5l, y5, feedback_x, y5)
    draw_line(sl, feedback_x, y5,
              feedback_x, y1)
    draw_arrow_right(sl, feedback_x, y1, LEFT)

    # Context management box
    CM_TOP = TOP + 2 * BH + 800000
    draw_box(sl, LEFT + 2 * (BW + GAP), CM_TOP,
             BW + 1500000, 900000,
             ["Context Management",
              "· Restart after 3-4 rounds",
              "· Recover via design doc",
              "· 1 module per context"],
             fill=WHITE, border=RGBColor(0xA5, 0xA5, 0xA5),
             font_size=Pt(11), align=PP_ALIGN.LEFT)

    # Constitution label
    draw_box(sl, LEFT, CM_TOP, 2 * BW + GAP, 900000,
             ["CLAUDE.md = Project Constitution",
              "98 KB design documentation",
              "Updated every iteration"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(12), align=PP_ALIGN.LEFT)


def slide_self_correct(sl):
    """Slide 7: Self-correction loop."""
    BW = 2200000
    BH = 650000
    GAP = 250000
    LEFT = 540000
    TOP = 1100000

    # Top row: Implement -> Write test -> Run test
    for i, text in enumerate(
            ["Implement", "Write test", "Run test"]):
        x = LEFT + i * (BW + GAP)
        draw_box(sl, x, TOP, BW, BH, [text],
                 font_size=Pt(14))
        if i < 2:
            draw_arrow_right(sl, x + BW,
                             TOP + BH // 2,
                             x + BW + GAP)

    # Branch from Run test
    run_cx = LEFT + 2 * (BW + GAP) + BW // 2
    branch_y = TOP + BH + 150000

    # FAIL branch (left)
    FAIL_X = LEFT + 1000000
    FAIL_Y = TOP + BH + 600000
    FAIL_W = 2200000
    FAIL_H = 900000
    draw_line(sl, run_cx, TOP + BH, run_cx, branch_y)
    draw_line(sl, run_cx, branch_y,
              FAIL_X + FAIL_W // 2, branch_y)
    draw_arrow_down(sl, FAIL_X + FAIL_W // 2, branch_y,
                    FAIL_Y)
    draw_label(sl, FAIL_X, branch_y - 250000,
               FAIL_W, 200000, "FAIL",
               font_size=Pt(12), color=ACCENT, bold=True)

    draw_box(sl, FAIL_X, FAIL_Y, FAIL_W, FAIL_H,
             ["Read spec / Read QEMU", "Analyze root cause"],
             fill=BG_LIGHT, font_size=Pt(12))

    # PASS branch (right)
    PASS_X = LEFT + 2 * (BW + GAP) + 400000
    PASS_Y = FAIL_Y
    PASS_W = 1800000
    PASS_H = 600000
    draw_line(sl, run_cx, branch_y,
              PASS_X + PASS_W // 2, branch_y)
    draw_arrow_down(sl, PASS_X + PASS_W // 2, branch_y,
                    PASS_Y)
    draw_label(sl, PASS_X, branch_y - 250000,
               PASS_W, 200000, "PASS",
               font_size=Pt(12), color=RGBColor(0x22, 0x8B, 0x22),
               bold=True)
    draw_box(sl, PASS_X, PASS_Y, PASS_W, PASS_H,
             ["Commit"],
             fill=RGBColor(0xE8, 0xF5, 0xE9),
             border=RGBColor(0x22, 0x8B, 0x22),
             font_size=Pt(14), bold=True)

    # Fix code box
    FIX_X = FAIL_X
    FIX_Y = FAIL_Y + FAIL_H + 300000
    FIX_W = FAIL_W
    FIX_H = 600000
    draw_box(sl, FIX_X, FIX_Y, FIX_W, FIX_H,
             ["Fix code"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(14), bold=True)
    draw_arrow_down(sl, FIX_X + FIX_W // 2, FAIL_Y + FAIL_H,
                    FIX_Y)

    # Loop back arrow from Fix to Run test
    loop_x = FIX_X + FIX_W
    loop_y = FIX_Y + FIX_H // 2
    run_r = LEFT + 2 * (BW + GAP) + BW
    draw_line(sl, loop_x, loop_y, run_r + 200000, loop_y)
    draw_line(sl, run_r + 200000, loop_y,
              run_r + 200000, TOP + BH // 2)
    draw_arrow_left(sl, run_r + 200000,
                    TOP + BH // 2, run_r)

    # Reference materials box
    REF_TOP = FIX_Y + FIX_H + 200000
    draw_box(sl, LEFT, REF_TOP, 10800000, 550000,
             ["Reference: RISC-V ISA manual · "
              "QEMU source · x86 SDM manual"],
             fill=WHITE, border=RGBColor(0xDD, 0xDD, 0xDD),
             font_size=Pt(11))


def slide_multi_agent(sl):
    """Slide 8: Dual-terminal collaboration."""
    BOX_W = 4500000
    BOX_H = 3200000
    GAP = 600000
    LEFT = 540000
    TOP = 1100000

    # Claude terminal
    draw_box(sl, LEFT, TOP, BOX_W, BOX_H,
             ["Claude Code", "",
              "· Architecture",
              "· Core logic",
              "· Code implementation",
              "· Design doc"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(13), align=PP_ALIGN.LEFT)
    draw_label(sl, LEFT, TOP - 50000, BOX_W, 300000,
               "Terminal A", font_size=Pt(11),
               color=ACCENT, bold=True)

    # Codex terminal
    RX = LEFT + BOX_W + GAP
    draw_box(sl, RX, TOP, BOX_W, BOX_H,
             ["Codex", "",
              "· Code Review",
              "· Test completion",
              "· Feedback"],
             fill=BG_LIGHT,
             border=RGBColor(0xA5, 0xA5, 0xA5),
             font_size=Pt(13), align=PP_ALIGN.LEFT)
    draw_label(sl, RX, TOP - 50000, BOX_W, 300000,
               "Terminal B", font_size=Pt(11),
               color=GRAY, bold=True)

    # Bidirectional arrow at exact edge midpoint
    center_y = TOP + BOX_H // 2
    conn = sl.shapes.add_connector(
        1, E(LEFT + BOX_W), E(center_y),
        E(RX), E(center_y))
    conn.line.color.rgb = ACCENT
    conn.line.width = Pt(2)
    ln = conn.line._ln
    for tag in ('a:headEnd', 'a:tailEnd'):
        end = ln.makeelement(qn(tag), {})
        end.set('type', 'triangle')
        end.set('w', 'med')
        end.set('len', 'med')
        ln.append(end)

    # Human orchestrates
    H_TOP = TOP + BOX_H + 200000
    HW = 3600000
    HH = 600000
    hx = LEFT + (2 * BOX_W + GAP - HW) // 2
    draw_box(sl, hx, H_TOP, HW, HH,
             ["Human orchestrates",
              "Different vendor models · Cross-verify"],
             fill=WHITE, border=ACCENT,
             font_size=Pt(12))
    # Lines from terminals to human box (inverted-Y at center)
    human_cx = hx + HW // 2
    join_y = (TOP + BOX_H + H_TOP) // 2
    cl_cx = LEFT + BOX_W // 2
    cx_cx = RX + BOX_W // 2
    draw_line(sl, cl_cx, TOP + BOX_H, cl_cx, join_y)
    draw_line(sl, cl_cx, join_y, human_cx, join_y)
    draw_line(sl, cx_cx, TOP + BOX_H, cx_cx, join_y)
    draw_line(sl, cx_cx, join_y, human_cx, join_y)
    draw_line(sl, human_cx, join_y, human_cx, H_TOP)


def slide_perf_story(sl):
    """Slide 9: Register mapping optimization story."""
    LEFT = 540000
    TOP = 1050000

    # Story box
    draw_box(sl, LEFT, TOP, 10800000, 1200000,
             ["Round 1: Optimization made it SLOWER",
              "",
              "Idea from Intel Houdini: "
              "guest regs → host regs direct mapping",
              "Expected: less context switch, 1:1 translation",
              "Result: dhrystone got SLOWER!"],
             fill=BG_LIGHT,
             border=RGBColor(0xA5, 0xA5, 0xA5),
             font_size=Pt(12), align=PP_ALIGN.LEFT)

    # Human insight box
    draw_box(sl, LEFT, TOP + 1400000, 10800000, 800000,
             ["Human's new idea: not enough x86 GPRs?",
              "→ Use SIMD regs (XMM) as GPR buffer",
              "→ Experimental, controlled by CLI flag"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(13), align=PP_ALIGN.LEFT)

    # Plan Mode flow: 3 boxes
    BW = 3200000
    BH = 700000
    GAP = 250000
    PM_TOP = TOP + 2600000
    for i, (t1, t2) in enumerate([
        ("Analyze", "avail XMM, conflict check"),
        ("Design", "mapping plan, key ops"),
        ("Output", "doc & change list"),
    ]):
        x = LEFT + i * (BW + GAP)
        draw_box(sl, x, PM_TOP, BW, BH, [t1, t2],
                 font_size=Pt(12))
        if i < 2:
            draw_arrow_right(sl, x + BW,
                             PM_TOP + BH // 2,
                             x + BW + GAP)

    draw_label(sl, LEFT, PM_TOP - 300000,
               10800000, 250000,
               "Agent enters Plan Mode (100K+ TOKENs)",
               font_size=Pt(12), color=GRAY,
               align=PP_ALIGN.LEFT)


def slide_perf_analysis(sl):
    """Slide 10: Performance analysis toolchain."""
    BW = 2400000
    BH = 750000
    GAP = 250000
    LEFT = 540000
    TOP = 1100000

    steps = [
        ("1. Benchmark", "dhrystone\ntime comparison"),
        ("2. Profiling", "perf sampling\nfind hotspot"),
        ("3. Code Quality", "IR translation\ninspection"),
        ("4. Optimize", "targeted fix\nre-benchmark"),
    ]
    for i, (t1, t2) in enumerate(steps):
        x = LEFT + i * (BW + GAP)
        lines = [t1] + t2.split("\n")
        fill = BG_ACCENT if i == 3 else BG_LIGHT
        brd = ACCENT if i == 3 else RGBColor(0xA5, 0xA5, 0xA5)
        draw_box(sl, x, TOP, BW, BH, lines,
                 fill=fill, border=brd, font_size=Pt(12))
        if i < 3:
            draw_arrow_right(sl, x + BW,
                             TOP + BH // 2,
                             x + BW + GAP)

    # Code quality checks box below
    CQ_TOP = TOP + BH + 400000
    draw_box(sl, LEFT, CQ_TOP, 10800000, 700000,
             ["Code Quality Checks: "
              "trace log (IR status) · "
              "host insn count (1 Guest → N Host) · "
              "register spill frequency"],
             fill=WHITE, border=RGBColor(0xDD, 0xDD, 0xDD),
             font_size=Pt(12), align=PP_ALIGN.LEFT)


def slide_difftest(sl):
    """Slide 12: Differential testing workflow."""
    BW = 2500000
    BH = 700000
    GAP = 1500000
    LEFT = 1500000
    TOP = 1100000

    # tcg-rs and QEMU boxes
    draw_box(sl, LEFT, TOP, BW, BH,
             ["tcg-rs", "exec instruction"],
             fill=BG_ACCENT, border=ACCENT, font_size=Pt(13))
    draw_box(sl, LEFT + BW + GAP, TOP, BW, BH,
             ["QEMU", "exec instruction"],
             fill=BG_LIGHT, font_size=Pt(13))

    # Reg snapshot boxes
    SNAP_Y = TOP + BH + 400000
    draw_box(sl, LEFT, SNAP_Y, BW, BH,
             ["Reg snapshot", "x0..x31"],
             fill=WHITE, border=ACCENT, font_size=Pt(12))
    draw_box(sl, LEFT + BW + GAP, SNAP_Y, BW, BH,
             ["Reg snapshot", "x0..x31"],
             fill=WHITE, border=RGBColor(0xA5, 0xA5, 0xA5),
             font_size=Pt(12))

    # Arrows down
    cx1 = LEFT + BW // 2
    cx2 = LEFT + BW + GAP + BW // 2
    draw_arrow_down(sl, cx1, TOP + BH, SNAP_Y)
    draw_arrow_down(sl, cx2, TOP + BH, SNAP_Y)

    # Compare box
    CMP_W = 2500000
    CMP_H = 700000
    cmp_x = LEFT + (2 * BW + GAP - CMP_W) // 2
    cmp_y = SNAP_Y + BH + 500000
    draw_box(sl, cmp_x, cmp_y, CMP_W, CMP_H,
             ["Compare", "Match?"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(14), bold=True)

    # Arrows from snapshots to compare
    mid_x = cmp_x + CMP_W // 2
    draw_line(sl, cx1, SNAP_Y + BH, cx1,
              SNAP_Y + BH + 200000)
    draw_line(sl, cx2, SNAP_Y + BH, cx2,
              SNAP_Y + BH + 200000)
    draw_line(sl, cx1, SNAP_Y + BH + 200000,
              mid_x, SNAP_Y + BH + 200000)
    draw_line(sl, cx2, SNAP_Y + BH + 200000,
              mid_x, SNAP_Y + BH + 200000)
    draw_arrow_down(sl, mid_x, SNAP_Y + BH + 200000,
                    cmp_y)

    # Stats labels
    stats_x = cmp_x + CMP_W + 300000
    draw_label(sl, stats_x, cmp_y,
               3000000, CMP_H,
               "816 tests\n35 diff-tests\n50% is test code",
               font_size=Pt(13), color=DARK,
               align=PP_ALIGN.LEFT)


def slide_expert_model(sl):
    """Slide 14: Expert experience x Model intelligence."""
    LEFT = 540000
    TOP = 1050000

    # Left column: Lever analogy
    LW = 5100000
    draw_box(sl, LEFT, TOP, LW, 600000,
             ["The Lever Analogy"],
             fill=BG_ACCENT, border=ACCENT,
             font_size=Pt(16), bold=True)

    # Lever visualization
    LEVER_Y = TOP + 750000
    LEVER_W = 4500000
    LEVER_H = 200000
    lever_x = LEFT + 300000
    # Lever bar
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            E(lever_x), E(LEVER_Y),
                            E(LEVER_W), E(LEVER_H))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()

    # Fulcrum triangle
    ful_x = lever_x + LEVER_W * 2 // 5
    ful_y = LEVER_Y + LEVER_H
    tri = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                              E(ful_x - 200000), E(ful_y),
                              E(400000), E(350000))
    tri.fill.solid()
    tri.fill.fore_color.rgb = DARK
    tri.line.fill.background()

    draw_label(sl, lever_x, LEVER_Y - 280000,
               1200000, 250000, "Effort",
               font_size=Pt(11), color=GRAY)
    draw_label(sl, lever_x + LEVER_W - 1200000,
               LEVER_Y - 280000, 1200000, 250000,
               "Outcome", font_size=Pt(11), color=GRAY)
    draw_label(sl, ful_x - 500000,
               ful_y + 350000, 1000000, 250000,
               "Fulcrum\n(Expert)", font_size=Pt(11),
               color=DARK, bold=True)

    # Key insight text
    INSIGHT_Y = LEVER_Y + 900000
    draw_box(sl, LEFT, INSIGHT_Y, LW, 1200000,
             ["Stronger model → longer lever",
              "",
              "→ Fulcrum covers LESS area",
              "  (no need for implementation detail)",
              "",
              "→ But position MORE critical",
              "  (direction judgment is key)"],
             fill=WHITE, border=RGBColor(0xDD, 0xDD, 0xDD),
             font_size=Pt(11), align=PP_ALIGN.LEFT)

    # Right column: Dynamic role table
    RX = LEFT + LW + 300000
    RW = 5400000
    add_table(sl, ["Phase", "Human Role", "Agent Autonomy"], [
        ["Early\n(backend)", "Deep: set path\npick strategy",
         "Low:\nfollow instructions"],
        ["Mid\n(frontend+IR)", "Define & verify\nreview design",
         "Med:\nplan + self-test"],
        ["Late\n(optimize)", "Direction only\naccept/reject",
         "High:\niterate autonomously"],
    ], left=RX, top=TOP, width=RW, row_height=480000)


def slide_three_practices(sl):
    """Slide 15: Three key practices."""
    BW = 3200000
    BH = 650000
    GAP = 200000
    LEFT = 540000

    practices = [
        ("#1: CLAUDE.md driven",
         [("Project charter", "98KB doc"),
          ("Restart after", "3-4 rounds"),
          ("New context", "Recover via doc")],
         BG_ACCENT, ACCENT),
        ("#2: Diff-test",
         [("Reference impl", "(QEMU)"),
          ("Agent self-test", "self-fix"),
          ("50% code", "is tests")],
         BG_LIGHT, RGBColor(0xA5, 0xA5, 0xA5)),
        ("#3: Multi-Agent",
         [("Claude", "codes"),
          ("Codex", "reviews"),
          ("Human routes", "No framework")],
         BG_LIGHT, RGBColor(0xA5, 0xA5, 0xA5)),
    ]

    for pi, (title, boxes, fill, brd) in enumerate(practices):
        y = 1050000 + pi * (BH + 500000)
        draw_label(sl, LEFT, y - 250000,
                   10800000, 230000, title,
                   font_size=Pt(13), color=ACCENT,
                   bold=True, align=PP_ALIGN.LEFT)
        for bi, (t1, t2) in enumerate(boxes):
            x = LEFT + bi * (BW + GAP)
            draw_box(sl, x, y, BW, BH, [t1, t2],
                     fill=fill if bi == 0 else WHITE,
                     border=brd, font_size=Pt(12))
            if bi < 2:
                draw_arrow_right(sl, x + BW,
                                 y + BH // 2,
                                 x + BW + GAP)


def slide_humanize(sl):
    """Slide 16: Humanize plugin."""
    LEFT = 540000
    TOP = 1050000
    LW = 3200000
    RW = 4800000
    GAP = 400000
    BH = 700000
    ROW_GAP = 150000
    RX = LEFT + LW + GAP + 800000  # space for arrow

    sections = [
        ("Manual dual terminal\nClaude + Codex",
         "RLCR auto loop\nClaude impl → Codex review\n→ feedback → iterate"),
        ("Doc-driven dev\nCLAUDE.md charter",
         "gen-plan output\nStructured plan (AC-X)\nTDD-style acceptance"),
        ("Agent self iteration\ntest & find bug",
         "Goal Tracker\nAlign review every 5 rounds\nPrevent drift"),
    ]

    labels = ["systemize", "structure", "safeguard"]

    for i, ((lt, rt), label) in enumerate(
            zip(sections, labels)):
        y = TOP + i * (BH + ROW_GAP)
        draw_box(sl, LEFT, y, LW, BH,
                 lt.split("\n"),
                 fill=BG_LIGHT,
                 border=RGBColor(0xA5, 0xA5, 0xA5),
                 font_size=Pt(11))
        draw_box(sl, RX, y, RW, BH,
                 rt.split("\n"),
                 fill=BG_ACCENT, border=ACCENT,
                 font_size=Pt(11), align=PP_ALIGN.LEFT)
        # Arrow with label
        ax1 = LEFT + LW
        ax2 = RX
        ay = y + BH // 2
        draw_arrow_right(sl, ax1, ay, ax2)
        draw_label(sl, ax1, ay - 200000,
                   ax2 - ax1, 180000, label,
                   font_size=Pt(9), color=ACCENT, bold=True)


# ============================================================
# Build presentation
# ============================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # SLIDE 1: Cover
    sl = add_blank_slide(prs)
    safe_add_picture(sl, IMG_BG, 0, 0,
                     SLIDE_W, SLIDE_H)
    safe_add_picture(sl, IMG_LOGO,
                     E(720000), E(540000),
                     E(1800000), E(750802))
    t = add_textbox(sl, 720000, 1620000, 7200000, 900000)
    set_run(t.text_frame.paragraphs[0], "tcg-rs",
            font_name=FONT_CODE, size=Pt(52),
            bold=True, color=ACCENT)
    t = add_textbox(sl, 720000, 2700000, 7920000, 1080000)
    set_run(t.text_frame.paragraphs[0],
            "AI Agent 驱动的二进制动态翻译系统",
            size=Pt(24), bold=True, color=DARK)
    set_run(t.text_frame.add_paragraph(), "构建与性能实践",
            size=Pt(24), bold=True, color=DARK)
    t = add_textbox(sl, 720000, 3960000, 7200000, 540000)
    set_run(t.text_frame.paragraphs[0],
            "太初元碁，高级软件工程师 刘超（泽文）",
            size=Pt(18), color=GRAY)
    t = add_textbox(sl, 720000, 4680000, 7200000, 432000)
    set_run(t.text_frame.paragraphs[0],
            "OS2ATC 2026 · AI 辅助编程分论坛",
            size=Pt(16), color=GRAY)

    # SLIDE 2: Background
    sl = add_blank_slide(prs)
    add_title(sl, "背景与挑战：为什么要重写 QEMU TCG")
    add_logo(sl)
    add_footer(sl, 2, TOTAL_SLIDES)
    add_code_block(sl, [
        "AI Chip Design Verification Flow",
        "",
        "RTL Sim ---> ISS Sim ---> FPGA ---> Tape-out",
        "(slowest)   (bottleneck!) (limited)  (highest cost)",
        "",
        "ISS needs high-perf Guest instruction execution",
        "Interpretation vs JIT: 10~50x performance gap",
    ], top=1008000, height=1800000, font_size=Pt(13))
    add_table(sl, ["维度", "QEMU TCG 的复杂度"], [
        ["历史", "25 年（2003 年诞生）"],
        ["核心代码", "~5 万行 C"],
        ["支持架构", "20+ Guest × 7 Host"],
        ["关键技术", "寄存器分配、指令编码、IR 优化、JIT 执行"],
        ["学习曲线", "x86 指令编码需 1-2 个月理解"],
    ], top=3000000, row_height=340000)
    add_highlight_box(sl,
        "核心问题：能否用 AI Agent 快速重塑这个引擎？",
        top=5200000, height=450000)

    # SLIDE 3: Pipeline
    sl = add_blank_slide(prs)
    add_title(sl, "动态二进制翻译：一图流")
    add_logo(sl)
    add_footer(sl, 3, TOTAL_SLIDES)
    slide_pipeline(sl)
    add_highlight_box(sl,
        "一句话解释：给你一本日语书，"
        "我把你想读的章节实时用中文朗读出来",
        top=5700000, height=500000)

    # SLIDE 4: Core thesis
    sl = add_blank_slide(prs)
    add_title(sl, "核心主张：人是产品经理，Agent 是工程团队")
    add_logo(sl)
    add_footer(sl, 4, TOTAL_SLIDES)
    slide_roles(sl)
    add_highlight_box(sl,
        "人只抓两件事：设计对不对 + 结果对不对",
        top=5500000, height=450000)

    # SLIDE 5: Agent decisions
    sl = add_blank_slide(prs)
    add_title(sl, "Agent 怎么做决策：后端先行 + Plan Mode")
    add_logo(sl)
    add_footer(sl, 5, TOTAL_SLIDES)
    slide_decisions(sl)
    add_quote(sl,
        "「它自己分析先实现哪部分指令，x86 的格式是什么样的，"
        "分几步走。你可以微调优先级，讨论出方案，"
        "然后它按流程一步步执行。」",
        top=3200000, height=600000)

    # SLIDE 6: Doc driven
    sl = add_blank_slide(prs)
    add_title(sl, "设计文档驱动 + 上下文管理")
    add_logo(sl)
    add_footer(sl, 6, TOTAL_SLIDES)
    slide_doc_driven(sl)

    # SLIDE 7: Self-iteration
    sl = add_blank_slide(prs)
    add_title(sl, "Agent 自我迭代：写测试、发现 Bug、自修复")
    add_logo(sl)
    add_footer(sl, 7, TOTAL_SLIDES)
    slide_self_correct(sl)
    add_quote(sl,
        "「大部分功能的错误，都是它自己写测试的时候发现的。"
        "我反而静态 review 代码，没有找出来太多。」",
        top=5800000, height=500000)

    # SLIDE 8: Multi-Agent
    sl = add_blank_slide(prs)
    add_title(sl, "多 Agent 协作：Claude 干活，Codex 审查")
    add_logo(sl)
    add_footer(sl, 8, TOTAL_SLIDES)
    slide_multi_agent(sl)
    add_highlight_box(sl,
        "为什么用两个？一是交叉验证互相卷；"
        "二是省钱——Claude 太贵了",
        top=5700000, height=450000)

    # SLIDE 9: Perf story
    sl = add_blank_slide(prs)
    add_title(sl, "性能优化实战：寄存器固定映射的故事")
    add_logo(sl)
    add_footer(sl, 9, TOTAL_SLIDES)
    slide_perf_story(sl)
    add_highlight_box(sl,
        "启示：优化方向要人类判断，Agent 负责工程落地和实验验证",
        top=5400000, height=450000)

    # SLIDE 10: Perf analysis
    sl = add_blank_slide(prs)
    add_title(sl, "Agent 怎么做性能分析")
    add_logo(sl)
    add_footer(sl, 10, TOTAL_SLIDES)
    slide_perf_analysis(sl)
    add_quote(sl,
        "「一整套流程几乎不用人参与。它自己就能找出来。"
        "这体现出通用智力强不强——"
        "它的语料库里能嗅到什么是好代码。」",
        top=3200000, height=600000)

    # SLIDE 11: Benchmark
    sl = add_blank_slide(prs)
    add_title(sl, "性能实测与诚实声明")
    add_logo(sl)
    add_footer(sl, 11, TOTAL_SLIDES)
    add_code_block(sl, [
        "dhrystone 200K iterations",
        "(RISC-V -> x86-64, linux-user)",
        "",
        "QEMU   ########################################  0.4~0.5s",
        "tcg-rs ################################          0.32s",
        "                                        ^",
        "                             observed 20%~30% faster",
    ], top=1008000, height=1600000, font_size=Pt(13))
    add_table(sl, ["优化类别", "估算贡献", "关键技术"], [
        ["执行循环", "~8-10%",
         "next_tb_hint + exit_target cache"],
        ["Guest 内存访问", "~8-10%",
         "No TLB, direct [R14+addr]"],
        ["数据结构 + 其他", "~8-13%",
         "Vec storage + zero-cost abstraction"],
    ], top=2800000, row_height=380000)
    add_code_block(sl, [
        "Honest Disclosure",
        ". Only ~30% features of original implemented",
        ". Less features -> shorter path -> expected",
        ". Data value: validates optimization direction",
    ], top=4600000, height=900000, font_size=Pt(13))

    # SLIDE 12: Difftest
    sl = add_blank_slide(prs)
    add_title(sl, "差分测试：不信 AI，信参考实现")
    add_logo(sl)
    add_footer(sl, 12, TOTAL_SLIDES)
    slide_difftest(sl)
    add_highlight_box(sl,
        "核心原则：QEMU 20 年打磨 >> 5 天 AI 生成 "
        "→ 以 QEMU 为 ground truth",
        top=5700000, height=450000)

    # SLIDE 13: AI boundary (left-right)
    sl = add_blank_slide(prs)
    add_title(sl, "AI 的边界与人类核心价值")
    add_logo(sl)
    add_footer(sl, 13, TOTAL_SLIDES)
    add_code_block(sl, [
        '"Peak at First Sight"',
        "",
        "AI code gives great first",
        "impression -- clean style,",
        "clear structure",
        "",
        "Looks like mid-level or",
        "near-senior engineer",
        "",
        "But deeper review reveals",
        "internal issues",
        "-> First look IS the",
        "   best look",
    ], left=540000, top=1008000,
       width=5100000, height=2500000, font_size=Pt(12))
    add_code_block(sl, [
        "Human core value =",
        "Define problems +",
        "Verify results",
        "",
        ". Knowledge >= AI to",
        "  collaborate effectively",
        '. AI forces learning --',
        '  or you get "fooled"',
        ". Production quality needs",
        "  human depth; AI is a",
        "  tool, not a replacement",
    ], left=540000, top=3700000,
       width=5100000, height=2400000, font_size=Pt(12))
    add_table(sl, ["AI 擅长", "AI 不擅长"], [
        ["C → Rust 对照翻译", "性能优化方向判断"],
        ["批量生成 184 条", "上下文变长后"],
        ["指令函数", "遗忘细节"],
        ["测试编写 +", "多线程竞态"],
        ["边界覆盖", "条件调试"],
        ["设计文档 +", "全局架构"],
        ["代码分析", "最优解"],
    ], left=5900000, top=1008000,
       width=5400000, row_height=380000)

    # SLIDE 14: Expert x Model (left-right)
    sl = add_blank_slide(prs)
    add_title(sl, "专家经验 × Agent 智力")
    add_logo(sl)
    add_footer(sl, 14, TOTAL_SLIDES)
    slide_expert_model(sl)
    add_highlight_box(sl,
        "模型越强，你不需要知道更多，但你需要判断得更准",
        top=6050000, height=400000)

    # SLIDE 15: Three practices
    sl = add_blank_slide(prs)
    add_title(sl, "三个关键实践")
    add_logo(sl)
    add_footer(sl, 15, TOTAL_SLIDES)
    slide_three_practices(sl)
    add_quote(sl,
        "这三个实践不是 tcg-rs 专属的，"
        "任何系统级 AI 辅助开发都可以借鉴",
        top=5500000, height=500000)

    # SLIDE 16: Humanize
    sl = add_blank_slide(prs)
    add_title(sl, "从手动到自动：Humanize 插件")
    add_logo(sl)
    add_footer(sl, 16, TOTAL_SLIDES)
    slide_humanize(sl)
    add_highlight_box(sl,
        "核心理念：迭代优于完美 · 独立审查防盲点 · "
        "多层防护防漂移",
        top=4800000, height=450000)
    add_quote(sl,
        "开源地址：github.com/humania-org/humanize",
        top=5500000, height=400000)

    # SLIDE 17: Data overview
    sl = add_blank_slide(prs)
    add_title(sl, "项目数据总览")
    add_logo(sl)
    add_footer(sl, 17, TOTAL_SLIDES)
    add_table(sl, ["指标", "数值", "说明"], [
        ["开发时间", "5 天", "全程 AI Agent 驱动"],
        ["总代码量", "~4.5 万行", "Rust（含生成代码和测试）"],
        ["测试数量", "816 个", "含差分测试 35 个，测试占比 50%"],
        ["Crate 数", "10 个", "分层解耦设计"],
        ["IR Opcodes", "158 个", "vs QEMU ~250（-40%）"],
        ["指令翻译", "184 条", "RV64IMAFDC 完整"],
        ["性能", "快 20%~30%", "vs QEMU（linux-user, dhrystone）"],
        ["花费", "~200 美元", "Claude + Codex"],
    ], top=1008000, row_height=400000)
    add_highlight_box(sl,
        "核心观点：AI 已具备系统级工程能力。"
        "关键不在于 AI 写代码有多强，"
        "而在于人类定义问题和验证结果的能力有多强。",
        top=4900000, height=600000)

    # SLIDE 18: Q&A
    sl = add_blank_slide(prs)
    safe_add_picture(sl, IMG_BG, 0, 0,
                     SLIDE_W, SLIDE_H)
    safe_add_picture(sl, IMG_LOGO,
                     E(720000), E(540000),
                     E(1800000), E(750802))
    t = add_textbox(sl, 720000, 2160000, 7200000, 900000)
    set_run(t.text_frame.paragraphs[0], "谢谢！",
            size=Pt(52), bold=True, color=ACCENT)
    t = add_textbox(sl, 720000, 3240000, 7200000, 540000)
    set_run(t.text_frame.paragraphs[0], "Q & A",
            size=Pt(28), bold=True, color=DARK)
    t = add_textbox(sl, 720000, 4140000, 7200000, 432000)
    set_run(t.text_frame.paragraphs[0],
            "github.com/patchfx/tcg-rs",
            size=Pt(16), color=GRAY)

    prs.save("os2atc-2026.pptx")
    print(f"Done! {len(prs.slides)} slides saved.")


build()
