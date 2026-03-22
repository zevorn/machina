"""
Generate os2atc-2026.pptx from os2atc-slides.md content.
Reuses images and style from the original PPTX.

All ASCII diagrams are built with helper functions to
guarantee box alignment. No hand-counted spaces.
"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -- constants --
SLIDE_W = 12192000
SLIDE_H = 6858000
ACCENT = RGBColor(0xE2, 0x6C, 0x0D)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
BG_ACCENT = RGBColor(0xFD, 0xF0, 0xE5)
FONT_CN = "微软雅黑"
FONT_CODE = "Consolas"
FOOTER_TEXT = "OS2ATC 2026 · AI 辅助编程分论坛"
TOTAL_SLIDES = 18


# ============================================================
# Box-drawing helpers — guarantees alignment
# ============================================================

def box(w, *lines):
    """Return (top, content_lines, bottom) for a box of
    inner width w.  Content lines are left-justified & padded.
    """
    top = "+" + "-" * w + "+"
    mid = ["|" + s.ljust(w) + "|" for s in lines]
    bot = top
    return (top, mid, bot)


def box_split(w, split_pos):
    """Bottom border with a split: +----+-----+ at split_pos."""
    left = "-" * split_pos
    right = "-" * (w - split_pos - 1)
    return "+" + left + "+" + right + "+"


def row(*cells, gap="  ", arrows=None):
    """Merge multiple box-tuples side by side.
    Returns list of strings.  arrows overrides gap for
    row index 1 (first content line).
    """
    # Unify to list of list-of-strings
    cols = []
    for c in cells:
        top, mid, bot = c
        cols.append([top] + mid + [bot])
    h = max(len(c) for c in cols)
    for c in cols:
        w = len(c[0])
        while len(c) < h:
            c.insert(-1, "|" + " " * (w - 2) + "|")
    result = []
    for r in range(h):
        parts = []
        for i, c in enumerate(cols):
            parts.append(c[r])
            if i < len(cols) - 1:
                if arrows and r == 1:
                    parts.append(arrows[i])
                else:
                    parts.append(gap)
        result.append("".join(parts))
    return result


def pad_lines(lines, width=None):
    """Pad all lines to the same width (max or given)."""
    if width is None:
        width = max(len(l) for l in lines)
    return [l.ljust(width) for l in lines]


# ============================================================
# Build diagram data
# ============================================================

def make_pipeline():
    """Slide 3: Translation Pipeline."""
    W = 8  # inner width
    b1 = box(W, "Guest   ", "ELF     ", "(RISC-V)")
    b2 = box(W, "Frontend", " Decode ", "        ")
    b3 = box(W, " TCG IR ", "  (IR)  ", "        ")
    b4_top = "+" + "-" * W + "+"
    b4_mid = [
        "|" + "Optimize".ljust(W) + "|",
        "|" + "ConstFld".ljust(W) + "|",
        "|" + "CopyProp".ljust(W) + "|",
    ]
    b4_bot = box_split(W, 4)  # +----+---+
    b4 = (b4_top, b4_mid, b4_bot)
    top_row = row(b1, b2, b3, b4, arrows=["->", "->", "->"])

    # Close box2 and box3 early (after 1 content line)
    # We do this by editing top_row lines
    bw = W + 2  # box width = 10
    g = 2       # gap
    # Row positions: box starts at 0, bw+g, 2*(bw+g), 3*(bw+g)
    # = 0, 12, 24, 36
    rw = 4 * bw + 3 * g  # row width = 46

    # Replace lines 3-4 (rows inside box2/3 that should
    # show closed bottoms)
    # Line 3 = box row index 3 (2nd content line)
    def set_cell(lines, line_idx, col_idx, text):
        """Replace a cell at (line_idx, col_idx) in lines."""
        start = col_idx * (bw + g)
        s = list(lines[line_idx])
        for i, ch in enumerate(text):
            s[start + i] = ch
        lines[line_idx] = "".join(s)

    # Pad top_row to rw
    top_row = pad_lines(top_row, rw)

    # Line 3 (index 3): close box2, box3 bottoms
    b2_bot = "+" + "-" * W + "+"
    b3_bot = "+" + "-" * W + "+"
    set_cell(top_row, 3, 1, b2_bot)
    set_cell(top_row, 3, 2, b3_bot)
    # Line 4 (index 4): close box1 bottom, box4 still open
    set_cell(top_row, 4, 0, "+" + "-" * W + "+")
    set_cell(top_row, 4, 1, " " * bw)
    set_cell(top_row, 4, 2, " " * bw)

    # Now the arrow-down line from box4 split
    split_col = 3 * (bw + g) + 4 + 1  # position of middle +
    # = 36 + 5 = 41

    lines = list(top_row)

    # Blank + vertical bar
    bar = " " * split_col + "|"
    lines.append(bar.ljust(rw))

    # Bottom row
    b5 = box(W, " Result ", "        ")
    b6 = box(W, "  JIT   ", "  Exec  ")
    b7 = box(W, "Backend ", " x86-64 ", "Codegen ")
    bot_row = row(b5, b6, b7, arrows=["<-", "<-"])
    bot_rw = 3 * bw + 2 * g  # 34

    for i, l in enumerate(bot_row):
        if i == 1:
            # First content line: add arrow to split
            arrow_len = split_col - len(l)
            l = l + "<" + "-" * (arrow_len - 2) + "+"
        lines.append(l.ljust(rw))

    # Close box 7 later (3 content lines)
    # The row helper already handled it; but we need to
    # remove the early close and keep Codegen line
    # Actually the row() already shows all content lines.
    # Let me just add the ^  and TB Cache below.

    # ^ pointing up from TB Cache to JIT
    jit_center = 1 * (bw + g) + bw // 2  # = 12 + 5 = 17
    backend_start = 2 * (bw + g)  # = 24
    caret_line = (" " * jit_center + "^"
                  + " " * (backend_start - jit_center - 1)
                  + "+" + "-" * W + "+")
    lines.append(caret_line.ljust(rw))

    tb_start = jit_center - W // 2 - 1  # center TB under ^
    tb = "+" + "-" * W + "+"
    tb_content = "|" + "TB Cache".ljust(W) + "|"
    lines.append(" " * tb_start + tb
                 + "  translate once, run many")
    lines.append(" " * tb_start + tb_content)
    lines.append(" " * tb_start + tb)

    return ["         Translation Pipeline", ""] + lines


def make_roles():
    """Slide 4: Role Division."""
    W = 14
    header = row(
        box(W, "    Human     ", " Product Mgr  ",
            "              "),
        box(W, "   Claude     ", " Architect +  ",
            "  Lead Dev    "),
        box(W, "    Codex     ", " Review +     ",
            "  Completion  "),
        gap=" ",
    )
    body = row(
        box(W, ". Define WHAT ", ". Define WHY  ",
            ". Verify      ", "  design doc  ",
            ". Verify test ", ". Perf        ",
            "  direction   "),
        box(W, ". Plan Mode   ", "  Decide &    ",
            "  Plan        ", ". Core logic  ",
            "  implement   ", ". Write       ",
            "  design doc  "),
        box(W, ". Code Review ", ". Batch test  ",
            ". Cross-      ", "  verify      ",
            ". Insn compl  ", "              ",
            "              "),
        gap=" ",
    )
    return (["               Role Division", ""]
            + header + body)


def make_decisions():
    """Slide 5: Plan Mode decision flow."""
    W = 10
    info = [
        "Human experience guides Agent's dev path",
        "+" + "-" * 48 + "+",
        "|" + " NOT: top-down -> split modules -> AI fills"
        .ljust(48) + "|",
        "|" + " YES: follow human's path: backend first"
        .ljust(48) + "|",
        "+" + "-" * 48 + "+",
        "",
        "Agent's decision process (Plan Mode)",
        "",
    ]
    b1 = box(W, "1.Analyze ", "x86 insn  ", "format    ")
    b2 = box(W, "2. Plan   ", "Prioritize", "Steps     ")
    # box3 with split bottom
    b3_top = "+" + "-" * W + "+"
    b3_mid = [
        "|" + "3.Discuss ".ljust(W) + "|",
        "|" + "Adjust    ".ljust(W) + "|",
        "|" + "Refine    ".ljust(W) + "|",
    ]
    b3_bot = box_split(W, 5)  # +-----+----+
    b3 = (b3_top, b3_mid, b3_bot)
    top_row = row(b1, b2, b3, gap=" ", arrows=[">", ">"])

    rw = 3 * (W + 2) + 2  # 38
    top_row = pad_lines(top_row, rw)

    split_pos = 2 * (W + 2 + 1) + 5 + 1  # box3 start + split
    # box3 starts at 2*(12+1) = 26, split at 26+6 = 32

    lines = info + top_row
    lines.append((" " * 32 + "|").ljust(rw))
    lines.append((" " * 32 + "v").ljust(rw))

    b4 = box(W, "4.Execute ", "Implement ")
    b4_top, b4_mid, b4_bot = b4
    # Position box4 starting at col 26 (same as box3)
    indent = 2 * (W + 2 + 1)  # 26
    lines.append((" " * indent + b4_top).ljust(rw))
    for m in b4_mid:
        lines.append((" " * indent + m).ljust(rw))
    lines.append((" " * indent + b4_bot).ljust(rw))

    return lines


def make_doc_driven():
    """Slide 6: Design doc iteration cycle."""
    W = 8
    bw = W + 2  # 10
    g = 2
    rw = 3 * bw + 2 * g  # 34

    b1 = box(W, "1.Define", "Require ")
    b2 = box(W, "2.Agent ", " Write  ", " design ", " doc    ")
    b3_top = "+" + "-" * W + "+"
    b3_mid = [
        "|" + "3.Agent ".ljust(W) + "|",
        "|" + "Generate".ljust(W) + "|",
        "|" + " code   ".ljust(W) + "|",
    ]
    b3_bot = box_split(W, 3)  # +---+----+
    b3 = (b3_top, b3_mid, b3_bot)
    top = row(b1, b2, b3, gap="  ", arrows=["->", "->"])
    top = pad_lines(top, rw)

    # Box1 closes at line 3 (index 3),
    # we need ^ underneath
    # Replace box1 bottom with +--------+, blank below
    # box1 is 3 lines (top + 2 content + bottom = 4 lines)
    # In top, line 0=top borders, 1=first content,
    # 2=second content, 3=bottom borders

    # Box2 has 4 content lines, box3 has 3 content lines
    # The row() helper already pads them.
    # Let me just check line count: max box height is box2
    # with 4 content = 6 lines total.

    # After the top section, add feedback loop
    # split_pos in box3: +---+----+ middle at
    # col 2*(bw+g) + 3 + 1 = 24 + 4 = 28
    split_col = 2 * (bw + g) + 3 + 1  # 28

    lines = [
        "CLAUDE.md = Project Constitution (98 KB)",
        "",
    ]

    # Manually build the cycle since boxes have
    # different lifecycles
    def L(s):
        return s.ljust(rw)

    # Row 0: top borders
    lines.append(L("+" + "-"*W + "+" + "  "
                    + "+" + "-"*W + "+" + "  "
                    + "+" + "-"*W + "+"))
    # Row 1: first content with arrows
    lines.append(L("|" + "1.Define".ljust(W) + "|"
                    + "->"
                    + "|" + "2.Agent ".ljust(W) + "|"
                    + "->"
                    + "|" + "3.Agent ".ljust(W) + "|"))
    # Row 2
    lines.append(L("|" + "Require ".ljust(W) + "|"
                    + "  "
                    + "|" + " Write  ".ljust(W) + "|"
                    + "  "
                    + "|" + "Generate".ljust(W) + "|"))
    # Row 3: box1 closes
    lines.append(L("+" + "-"*W + "+"
                    + "  "
                    + "|" + " design ".ljust(W) + "|"
                    + "  "
                    + "|" + " code   ".ljust(W) + "|"))
    # Row 4: ^ under box1, box2 continues, box3 closes
    col1_center = bw // 2  # 5
    gap_to_b2 = (bw + g) - col1_center - 1  # 12-5-1=6
    lines.append(L(" " * col1_center + "^"
                    + " " * gap_to_b2
                    + "|" + " doc    ".ljust(W) + "|"
                    + "  "
                    + box_split(W, 3)))
    # Row 5: | under ^, box2 closes, | from split
    box2_end_pos = (bw + g) + bw  # 22
    gap_to_split = split_col - box2_end_pos  # 6
    lines.append(L(" " * col1_center + "|"
                    + " " * gap_to_b2
                    + "+" + "-"*W + "+"
                    + " " * gap_to_split
                    + "|"))
    # Row 6: | under ^, box5 top, box4 top with v
    b5_start = bw + g  # 12
    b4_start = 2 * (bw + g)  # 24
    b4_top_split = "+" + "-"*3 + "v" + "-"*4 + "+"
    lines.append(L(" " * col1_center + "|"
                    + " " * (b5_start - col1_center - 1)
                    + "+" + "-"*W + "+"
                    + "  "
                    + b4_top_split))
    # Row 7: feedback arrow + box5 content + box4 content
    arrow = "+" + "-" * (b5_start - col1_center - 1)
    lines.append(L(" " * col1_center + arrow
                    + "|" + "5.Update".ljust(W) + "|"
                    + "<-"
                    + "|" + "4. Diff ".ljust(W) + "|"))
    # Row 8-9: box5 & box4 content
    lines.append(L(" " * b5_start
                    + "|" + " design ".ljust(W) + "|"
                    + "  "
                    + "|" + " test   ".ljust(W) + "|"))
    lines.append(L(" " * b5_start
                    + "|" + " doc    ".ljust(W) + "|"
                    + "  "
                    + "|" + " verify ".ljust(W) + "|"))
    # Row 10: bottom borders
    lines.append(L(" " * b5_start
                    + "+" + "-"*W + "+"
                    + "  "
                    + "+" + "-"*W + "+"))
    return lines


def make_self_correct():
    """Slide 7: Self-correction loop."""
    W = 10
    b1 = box(W, " Implement", "          ")
    b2 = box(W, "Write test", "          ")
    b3 = box(W, " Run test ", "          ")
    top = row(b1, b2, b3, gap="  ", arrows=["->", "->"])
    bw = W + 2  # 12
    g = 2
    rw = 3 * bw + 2 * g  # 40
    top = pad_lines(top, rw)

    # Split at bottom of box3
    split_col = 2 * (bw + g) + bw // 2  # 28+6=34
    lines = ["Agent's Self-Correction Loop", ""]
    lines += top

    lines.append((" " * split_col + "|").ljust(rw))

    # Branch: position boxes with proper gap
    bf_indent = 15   # fail box start
    bp_indent = 29   # pass box start
    gap_between = bp_indent - bf_indent - bw  # 2

    # Draw branch lines from split_col
    fail_bar = bf_indent + bw // 2   # 21
    pass_bar = bp_indent + bw // 2   # 35
    lines.append((" " * fail_bar + "+"
                  + "-" * (pass_bar - fail_bar - 1)
                  + "+").ljust(rw))
    lines.append((" " * (fail_bar - 4) + "FAIL"
                  + " " * 1 + "|"
                  + " " * (pass_bar - fail_bar - 1)
                  + "|"
                  + " " + "PASS").ljust(rw))
    lines.append((" " * fail_bar + "v"
                  + " " * (pass_bar - fail_bar - 1)
                  + "v").ljust(rw))

    # Two boxes side by side
    bf = box(W, " Read spec", " Read QEMU", " Analyze  ")
    bp = box(W, "  Commit  ")
    bf_top, bf_mid, bf_bot = bf
    bp_top, bp_mid, bp_bot = bp

    lines.append((" " * bf_indent + bf_top
                  + " " * gap_between + bp_top).ljust(rw))
    for i in range(max(len(bf_mid), len(bp_mid))):
        fl = bf_mid[i] if i < len(bf_mid) else " " * bw
        pl = bp_mid[i] if i < len(bp_mid) else " " * bw
        lines.append((" " * bf_indent + fl
                      + " " * gap_between + pl).ljust(rw))
    lines.append((" " * bf_indent + bf_bot
                  + " " * gap_between + bp_bot).ljust(rw))

    # Arrow down from fix
    fix_center = bf_indent + bw // 2
    lines.append((" " * fix_center + "|").ljust(rw))
    lines.append((" " * fix_center + "v").ljust(rw))

    bfix = box(W, " Fix code ")
    fx_top, fx_mid, fx_bot = bfix
    lines.append((" " * bf_indent + fx_top).ljust(rw))
    lines.append((" " * bf_indent + fx_mid[0]
                  + "-> (to Run test)").ljust(rw))
    lines.append((" " * bf_indent + fx_bot).ljust(rw))

    return lines


def make_multi_agent():
    """Slide 8: Dual-terminal collaboration."""
    W = 18
    bw = W + 2  # 20
    g = 5
    rw = 2 * bw + g  # 45

    def L(s):
        return s.ljust(rw)

    b = "+" + "-" * W + "+"
    lines = [
        "Dual-Terminal Collaboration",
        "",
        L(b + " " * g + b),
        L("|" + " " * W + "|" + " " * g
          + "|" + " " * W + "|"),
        L("|" + "   Claude Code    ".ljust(W) + "|"
          + " " * g
          + "|" + "     Codex        ".ljust(W) + "|"),
        L("|" + " " * W + "|" + " " * g
          + "|" + " " * W + "|"),
        L("|" + " . Architecture   ".ljust(W) + "|"
          + " " * g
          + "|" + " . Code Review    ".ljust(W) + "|"),
        L("|" + " . Core logic     ".ljust(W) + "|"
          + " --> "
          + "|" + " . Test compl     ".ljust(W) + "|"),
        L("|" + " . Code impl      ".ljust(W) + "|"
          + " " * g
          + "|" + " . Feedback       ".ljust(W) + "|"),
        L("|" + " . Design doc     ".ljust(W) + "|"
          + " <-- "
          + "|" + " " * W + "|"),
        L("|" + " " * W + "|" + " " * g
          + "|" + " " * W + "|"),
        L(b + " " * g + b),
    ]

    # Below: human orchestrates
    c1 = bw // 2  # 10
    c2 = bw + g + bw // 2  # 35
    lines.append(L(" " * c1 + "^"
                    + " " * (c2 - c1 - 1) + "^"))
    lines.append(L(" " * c1 + "|"
                    + "  Human orchestrates  "
                    + "|"))
    mid = (c1 + c2) // 2
    lines.append(L(" " * c1 + "+"
                    + "-" * (c2 - c1 - 1) + "+"))
    lines.append(L(" " * mid + "|"))

    bx_w = 20
    bx_start = mid - bx_w // 2
    bx = "+" + "-" * (bx_w - 2) + "+"
    lines.append(L(" " * bx_start + bx))
    lines.append(L(" " * bx_start
                    + "|" + "Different vendors ".ljust(bx_w-2)
                    + "|"))
    lines.append(L(" " * bx_start
                    + "|" + "Top models       ".ljust(bx_w-2)
                    + "|"))
    lines.append(L(" " * bx_start
                    + "|" + "Cross-verify     ".ljust(bx_w-2)
                    + "|"))
    lines.append(L(" " * bx_start + bx))

    return lines


def make_perf_story():
    """Slide 9: Register mapping story."""
    W = 12
    lines = [
        "Round 1: Optimization made it SLOWER",
        "",
        "Idea from Intel Houdini:",
        "  guest regs -> host regs direct mapping",
        "Expected: less context switch, 1:1 translation",
        "Result:   dhrystone got SLOWER!",
        "",
    ]
    # info box
    iw = 44
    ib = "+" + "-" * iw + "+"
    lines.append(ib)
    lines.append("|" + " Human's new idea: not enough x86 GPRs?"
                 .ljust(iw) + "|")
    lines.append("|" + " -> Use SIMD regs (XMM) as GPR buffer"
                 .ljust(iw) + "|")
    lines.append("|" + " -> Experimental, controlled by CLI flag"
                 .ljust(iw) + "|")
    lines.append(ib)
    lines.append("")
    lines.append("Agent enters Plan Mode (100K+ TOKENs)")
    lines.append("")

    b1 = box(W, "Analyze     ", "avail XMM   ", "Conflict chk")
    b2 = box(W, "Design      ", "mapping     ", "Key ops     ")
    b3 = box(W, "Output      ", "doc &       ", "change list ")
    r = row(b1, b2, b3, gap="  ", arrows=["->", "->"])
    lines += r
    return lines


def make_perf_analysis():
    """Slide 10: Performance analysis toolchain."""
    W = 12
    bw = W + 2  # 14
    g = 2
    rw = 3 * bw + 2 * g  # 46

    b1 = box(W, "1.Benchmark ", "dhrystone   ", "time cmp    ")
    b2 = box(W, "2.Profiling ", "perf sample ", "find hotspot")
    b3_top = "+" + "-" * W + "+"
    b3_mid = [
        "|" + "3.Code      ".ljust(W) + "|",
        "|" + "quality     ".ljust(W) + "|",
        "|" + "inspection  ".ljust(W) + "|",
    ]
    b3_bot = box_split(W, 5)
    b3 = (b3_top, b3_mid, b3_bot)
    top = row(b1, b2, b3, gap="  ", arrows=["->", "->"])
    top = pad_lines(top, rw)

    split_col = 2 * (bw + g) + 5 + 1  # 32+6=38
    lines = ["Agent's Performance Analysis Toolchain", ""]
    lines += top
    lines.append((" " * split_col + "|").ljust(rw))
    lines.append((" " * split_col + "v").ljust(rw))

    b4 = box(W, "4.Targeted  ", "optimize    ", "re-benchmark")
    indent = 2 * (bw + g)  # 32
    b4t, b4m, b4b = b4
    lines.append((" " * indent + b4t).ljust(rw))
    for m in b4m:
        lines.append((" " * indent + m).ljust(rw))
    lines.append((" " * indent + b4b).ljust(rw))
    return lines


def make_difftest():
    """Slide 12: Differential testing."""
    W = 10
    bw = W + 2  # 12
    g = 5
    rw = 2 * bw + g  # 29

    def L(s):
        return s.ljust(rw)

    b = "+" + "-" * W + "+"
    sb = box_split(W, 4)  # +----+-----+

    lines = ["Differential Testing Workflow", ""]
    lines.append(L(b + " " * g + b))
    lines.append(L("|" + "  tcg-rs  ".ljust(W) + "|"
                    + " " * g
                    + "|" + "   QEMU   ".ljust(W) + "|"))
    lines.append(L("|" + " exec insn".ljust(W) + "|"
                    + " " * g
                    + "|" + " exec insn".ljust(W) + "|"))
    lines.append(L(sb + " " * g + sb))

    c1 = 4 + 1  # split_col of box1
    c2 = bw + g + 4 + 1  # split_col of box2
    lines.append(L(" " * c1 + "|"
                    + " " * (c2 - c1 - 1) + "|"))
    lines.append(L(" " * c1 + "v"
                    + " " * (c2 - c1 - 1) + "v"))

    lines.append(L(b + " " * g + b))
    lines.append(L("|" + "Reg snap  ".ljust(W) + "|"
                    + " " * g
                    + "|" + "Reg snap  ".ljust(W) + "|"))
    lines.append(L("|" + " x0..x31  ".ljust(W) + "|"
                    + " " * g
                    + "|" + " x0..x31  ".ljust(W) + "|"))
    lines.append(L(sb + " " * g + sb))

    lines.append(L(" " * c1 + "|"
                    + " " * (c2 - c1 - 1) + "|"))
    mid = (c1 + c2) // 2
    lines.append(L(" " * c1 + "+"
                    + "-" * (c2 - c1 - 1) + "+"))
    lines.append(L(" " * mid + "|"))
    lines.append(L(" " * mid + "v"))

    cw = 12
    ci = mid - cw // 2
    cb = "+" + "-" * (cw - 2) + "+"
    lines.append(L(" " * ci + cb) + "  816 tests")
    lines.append(L(" " * ci
                    + "|" + " Compare ".ljust(cw-2) + "|")
                 + "  35 diff-tests")
    lines.append(L(" " * ci
                    + "|" + " Match?  ".ljust(cw-2) + "|")
                 + "  50% is test code")
    lines.append(L(" " * ci + cb))
    lines.append(L(" " * ci + " PASS | FAIL -> locate bug"))
    return lines


def make_three_practices():
    """Slide 14: Three practices."""
    W = 8
    lines = []
    lines.append("#1: CLAUDE.md -- design doc as"
                 " persistent memory")
    r1 = row(
        box(W, "Project ", "charter ", "98KB doc"),
        box(W, "Restart ", "after   ", "3-4 rnds"),
        box(W, "New ctx ", "Recover ", "via doc "),
        gap="  ", arrows=["->", "->"])
    lines += r1
    lines.append("-" * 34)
    lines.append("#2: Diff-test -- trust reference,"
                 " not AI")
    r2 = row(
        box(W, "Ref impl", "(QEMU)  ", "        "),
        box(W, "Agent   ", "selftest", "selffix "),
        box(W, "50% code", "is tests", "        "),
        gap="  ", arrows=["->", "->"])
    lines += r2
    lines.append("-" * 34)
    lines.append("#3: Multi-Agent -- different models"
                 " compete")
    r3 = row(
        box(W, " Claude ", " codes  "),
        box(W, " Codex  ", " reviews"),
        box(W, " Human  ", " routes ", "No frmwk"),
        gap="  ", arrows=["<>", "->"])
    lines += r3
    return lines


def make_humanize():
    """Slide 15: Humanize plugin."""
    LW = 12  # left box inner width
    RW = 18  # right box inner width
    g = 6    # gap between columns
    lbw = LW + 2  # 14
    rbw = RW + 2  # 20
    rw = lbw + g + rbw  # 40

    def L(s):
        return s.ljust(rw)

    lb = "+" + "-" * LW + "+"
    rb = "+" + "-" * RW + "+"

    def lr(lt, rt, arrow=False):
        mid = "===>> " if arrow else " " * g
        return L(lt + mid + rt)

    lines = [
        "tcg-rs Practice       Humanize Systematic",
        "",
    ]
    lines.append(lr(lb, rb))
    lines.append(lr(
        "|" + "Manual dual  ".ljust(LW) + "|",
        "|" + "RLCR auto loop    ".ljust(RW) + "|"))
    lines.append(lr(
        "|" + "terminal    ".ljust(LW) + "|",
        "|" + "Claude implement  ".ljust(RW) + "|",
        arrow=True))
    lines.append(lr(
        "|" + "Claude+Codex".ljust(LW) + "|",
        "|" + "-> Codex review   ".ljust(RW) + "|"))
    lines.append(lr(
        lb,
        "|" + "-> feedback iter  ".ljust(RW) + "|"))
    lines.append(lr(" " * lbw, rb))

    lines.append(lr(lb,
        "|" + "gen-plan output   ".ljust(RW) + "|"))
    lines.append(lr(
        "|" + "Doc-driven  ".ljust(LW) + "|",
        "|" + "structured plan   ".ljust(RW) + "|"))
    lines.append(lr(
        "|" + "CLAUDE.md   ".ljust(LW) + "|",
        "|" + "(AC-X criteria)   ".ljust(RW) + "|",
        arrow=True))
    lines.append(lr(
        "|" + "charter     ".ljust(LW) + "|",
        "|" + "TDD-style accept  ".ljust(RW) + "|"))
    lines.append(lr(lb, rb))

    lines.append(lr(lb,
        "|" + "Goal Tracker      ".ljust(RW) + "|"))
    lines.append(lr(
        "|" + "Agent self  ".ljust(LW) + "|",
        "|" + "align review      ".ljust(RW) + "|"))
    lines.append(lr(
        "|" + "iteration   ".ljust(LW) + "|",
        "|" + "every 5 rounds    ".ljust(RW) + "|",
        arrow=True))
    lines.append(lr(
        "|" + "test & fix  ".ljust(LW) + "|",
        "|" + "prevent drift     ".ljust(RW) + "|"))
    lines.append(lr(lb, rb))

    return lines


# ============================================================
# PPTX building helpers
# ============================================================

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_run(p, text, font_name=FONT_CN, size=Pt(14),
            bold=False, color=DARK, italic=False):
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return r


def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def add_accent_bar(slide):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Emu(324000),
        Emu(108000), Emu(504000))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def add_title(slide, text):
    add_accent_bar(slide)
    t = tb(slide, Emu(252000), Emu(270000),
           Emu(9000000), Emu(576000))
    set_run(t.text_frame.paragraphs[0], text,
            size=Pt(26), color=ACCENT)


def add_footer(slide, n, total):
    t = tb(slide, Emu(540000), Emu(6480000),
           Emu(6480000), Emu(252000))
    set_run(t.text_frame.paragraphs[0], FOOTER_TEXT,
            size=Pt(11), color=GRAY)
    t2 = tb(slide, Emu(10800000), Emu(6480000),
            Emu(1080000), Emu(252000))
    t2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_run(t2.text_frame.paragraphs[0], f"{n}/{total}",
            size=Pt(11), color=GRAY)


def add_logo(slide):
    slide.shapes.add_picture(
        "/tmp/pptx_logo2.png",
        Emu(10260000), Emu(108000),
        Emu(1368000), Emu(570610))


def add_code_block(slide, lines, left=540000, top=1008000,
                   width=10800000, height=5200000,
                   font_size=Pt(12)):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top),
        Emu(width), Emu(height))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    s.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    s.line.width = Pt(1)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(180000)
    tf.margin_top = Emu(144000)
    tf.margin_right = Emu(180000)
    tf.margin_bottom = Emu(144000)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        p.space_before = Pt(0)
        set_run(p, line, font_name=FONT_CODE,
                size=font_size, color=DARK)


def add_table(slide, headers, rows, left=540000, top=1008000,
              width=10800000, row_height=360000):
    tr = 1 + len(rows)
    s = slide.shapes.add_table(
        tr, len(headers),
        Emu(left), Emu(top),
        Emu(width), Emu(tr * row_height))
    t = s.table
    cw = width // len(headers)
    for i in range(len(headers)):
        t.columns[i].width = Emu(cw)
    for ci, h in enumerate(headers):
        c = t.cell(0, ci)
        c.text = ""
        set_run(c.text_frame.paragraphs[0], h,
                size=Pt(14), bold=True, color=WHITE)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
    for ri, r in enumerate(rows):
        for ci, v in enumerate(r):
            c = t.cell(ri + 1, ci)
            c.text = ""
            set_run(c.text_frame.paragraphs[0], v,
                    size=Pt(13), color=DARK)
            c.fill.solid()
            c.fill.fore_color.rgb = (
                BG_LIGHT if ri % 2 == 0 else WHITE)


def add_highlight_box(slide, text, left=540000, top=5400000,
                      width=10800000, height=540000):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height))
    s.fill.solid()
    s.fill.fore_color.rgb = BG_ACCENT
    s.line.color.rgb = ACCENT
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(180000)
    tf.margin_top = Emu(72000)
    set_run(tf.paragraphs[0], text,
            size=Pt(14), bold=True, color=ACCENT)


def add_quote(slide, text, left=540000, top=5400000,
              width=10800000, height=500000):
    b = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(left), Emu(top), Emu(54000), Emu(height))
    b.fill.solid()
    b.fill.fore_color.rgb = ACCENT
    b.line.fill.background()
    t = tb(slide, Emu(left + 108000), Emu(top),
           Emu(width - 108000), Emu(height))
    t.text_frame.word_wrap = True
    set_run(t.text_frame.paragraphs[0], text,
            size=Pt(13), italic=True, color=GRAY)


def make_expert_model():
    """Slide 14: Expert experience x Model intelligence."""
    # Lever analogy diagram
    # lever_w = 50
    lines = [
        "The Lever Analogy",
        "",
        "Model intelligence = lever length",
        "Expert experience  = fulcrum position",
        "",
        "          effort                     outcome",
        "            |                           |",
        "            v                           v",
        "  +---------+---------------------------+---+",
        "  |/////////|                           |   |",
        "  +---------+---------------------------+---+",
        "                   ^",
        "                   |",
        "              +---------+",
        "              | Fulcrum |",
        "              | (Expert)|",
        "              +---------+",
        "",
        "Stronger model -> longer lever",
        "  -> fulcrum covers LESS area (no need for detail)",
        "  -> but fulcrum position MORE critical (direction)",
    ]

    return lines


def make_dynamic_roles():
    """Dynamic role table for slide 14."""
    W = 14
    bw = W + 2  # 16
    g = 1
    rw = 3 * bw + 2 * g  # 50

    def L(s):
        return s.ljust(rw)

    hdr = "+" + "-" * W + "+"
    lines = [
        "Dynamic Role Shift Across Project Phases",
        "",
        L(hdr + " " * g + hdr + " " * g + hdr),
        L("|" + "    Phase     ".ljust(W) + "|"
          + " " * g
          + "|" + " Human's Role ".ljust(W) + "|"
          + " " * g
          + "|" + "Agent Autonomy".ljust(W) + "|"),
        L(hdr + " " * g + hdr + " " * g + hdr),
        L("|" + "Early         ".ljust(W) + "|"
          + " " * g
          + "|" + "Deep: set path".ljust(W) + "|"
          + " " * g
          + "|" + "Low: follow   ".ljust(W) + "|"),
        L("|" + "(backend)     ".ljust(W) + "|"
          + " " * g
          + "|" + "pick strategy ".ljust(W) + "|"
          + " " * g
          + "|" + "instructions  ".ljust(W) + "|"),
        L(hdr + " " * g + hdr + " " * g + hdr),
        L("|" + "Mid           ".ljust(W) + "|"
          + " " * g
          + "|" + "Define & verif".ljust(W) + "|"
          + " " * g
          + "|" + "Med: plan +   ".ljust(W) + "|"),
        L("|" + "(frontend+IR) ".ljust(W) + "|"
          + " " * g
          + "|" + "review design ".ljust(W) + "|"
          + " " * g
          + "|" + "self-test     ".ljust(W) + "|"),
        L(hdr + " " * g + hdr + " " * g + hdr),
        L("|" + "Late          ".ljust(W) + "|"
          + " " * g
          + "|" + "Direction only".ljust(W) + "|"
          + " " * g
          + "|" + "High: iterate ".ljust(W) + "|"),
        L("|" + "(optimize)    ".ljust(W) + "|"
          + " " * g
          + "|" + "accept/reject ".ljust(W) + "|"
          + " " * g
          + "|" + "autonomously  ".ljust(W) + "|"),
        L(hdr + " " * g + hdr + " " * g + hdr),
    ]
    return lines


# ============================================================
# Build presentation
# ============================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # SLIDE 1: Cover
    sl = add_blank_slide(prs)
    sl.shapes.add_picture("/tmp/pptx_bg.png", 0, 0,
                          SLIDE_W, SLIDE_H)
    sl.shapes.add_picture("/tmp/pptx_logo.png",
                          Emu(720000), Emu(540000),
                          Emu(1800000), Emu(750802))
    t = tb(sl, Emu(720000), Emu(1620000),
           Emu(7200000), Emu(900000))
    set_run(t.text_frame.paragraphs[0], "tcg-rs",
            font_name=FONT_CODE, size=Pt(52),
            bold=True, color=ACCENT)
    t = tb(sl, Emu(720000), Emu(2700000),
           Emu(7920000), Emu(1080000))
    set_run(t.text_frame.paragraphs[0],
            "AI Agent 驱动的二进制动态翻译系统",
            size=Pt(24), bold=True, color=DARK)
    set_run(t.text_frame.add_paragraph(), "构建与性能实践",
            size=Pt(24), bold=True, color=DARK)
    t = tb(sl, Emu(720000), Emu(3960000),
           Emu(7200000), Emu(540000))
    set_run(t.text_frame.paragraphs[0],
            "太初元碁，高级软件工程师 刘超（泽文）",
            size=Pt(18), color=GRAY)
    t = tb(sl, Emu(720000), Emu(4680000),
           Emu(7200000), Emu(432000))
    set_run(t.text_frame.paragraphs[0],
            "OS2ATC 2026 · AI 辅助编程分论坛",
            size=Pt(16), color=GRAY)

    # SLIDE 2: Background
    sl = add_blank_slide(prs)
    add_title(sl, "背景与挑战：为什么要重写 QEMU TCG")
    add_logo(sl)
    add_footer(sl, 2, TOTAL_SLIDES)
    add_code_block(sl, [
        "AI Chip Design Verification Flow",
        "",
        "RTL Sim ---> ISS Sim ---> FPGA ---> Tape-out",
        "(slowest)   (bottleneck!) (limited)  (highest cost)",
        "",
        "ISS needs high-perf Guest instruction execution",
        "Interpretation vs JIT: 10~50x performance gap",
    ], top=1008000, height=1800000, font_size=Pt(13))
    add_table(sl, ["维度", "QEMU TCG 的复杂度"], [
        ["历史", "25 年（2003 年诞生）"],
        ["核心代码", "~5 万行 C"],
        ["支持架构", "20+ Guest × 7 Host"],
        ["关键技术", "寄存器分配、指令编码、IR 优化、JIT 执行"],
        ["学习曲线", "x86 指令编码需 1-2 个月理解"],
    ], top=3000000, row_height=340000)
    add_highlight_box(sl,
        "核心问题：能否用 AI Agent 快速重塑这个引擎？",
        top=5200000, height=450000)

    # SLIDE 3: Pipeline
    sl = add_blank_slide(prs)
    add_title(sl, "动态二进制翻译：一图流")
    add_logo(sl)
    add_footer(sl, 3, TOTAL_SLIDES)
    add_code_block(sl, make_pipeline(), top=1008000,
                   height=3800000, font_size=Pt(11))
    add_highlight_box(sl,
        "一句话解释：给你一本日语书，"
        "我把你想读的章节实时用中文朗读出来",
        top=5100000, height=500000)

    # SLIDE 4: Core thesis
    sl = add_blank_slide(prs)
    add_title(sl, "核心主张：人是产品经理，Agent 是工程团队")
    add_logo(sl)
    add_footer(sl, 4, TOTAL_SLIDES)
    add_code_block(sl, make_roles(), top=1008000,
                   height=3800000, font_size=Pt(10))
    add_highlight_box(sl,
        "人只抓两件事：设计对不对 + 结果对不对",
        top=5100000, height=450000)

    # SLIDE 5: Agent decisions
    sl = add_blank_slide(prs)
    add_title(sl, "Agent 怎么做决策：后端先行 + Plan Mode")
    add_logo(sl)
    add_footer(sl, 5, TOTAL_SLIDES)
    add_code_block(sl, make_decisions(), top=1008000,
                   height=4200000, font_size=Pt(10))
    add_quote(sl,
        "「它自己分析先实现哪部分指令，x86 的格式是什么样的，"
        "分几步走。你可以微调优先级，讨论出方案，"
        "然后它按流程一步步执行。」",
        top=5400000, height=600000)

    # SLIDE 6: Doc driven
    sl = add_blank_slide(prs)
    add_title(sl, "设计文档驱动 + 上下文管理")
    add_logo(sl)
    add_footer(sl, 6, TOTAL_SLIDES)
    add_code_block(sl, make_doc_driven(), top=1008000,
                   height=2700000, font_size=Pt(11))
    add_code_block(sl, [
        "Context Management Strategy",
        "",
        ". Terminate after 3-4 rounds -> new context",
        ". Recover via design doc -> 1 module per ctx",
        ". Effect: lower TOKEN cost + higher accuracy",
    ], top=3900000, height=1100000, font_size=Pt(12))
    add_quote(sl,
        "「类似芯片设计：所有实现全部对齐同一份文档。"
        "AI Agent 也是一样的思路。」",
        top=5200000, height=500000)

    # SLIDE 7: Self-iteration
    sl = add_blank_slide(prs)
    add_title(sl, "Agent 自我迭代：写测试、发现 Bug、自修复")
    add_logo(sl)
    add_footer(sl, 7, TOTAL_SLIDES)
    add_code_block(sl, make_self_correct(), top=1008000,
                   height=3600000, font_size=Pt(10))
    add_code_block(sl, [
        "Reference Materials",
        ". RISC-V ISA manual  -> instruction semantics",
        ". QEMU source code   -> expected behavior",
        ". x86 SDM manual     -> backend encoding",
    ], top=4700000, height=900000, font_size=Pt(12))
    add_quote(sl,
        "「大部分功能的错误，都是它自己写测试的时候发现的。"
        "我反而静态 review 代码，没有找出来太多。」",
        top=5750000, height=500000)

    # SLIDE 8: Multi-Agent
    sl = add_blank_slide(prs)
    add_title(sl, "多 Agent 协作：Claude 干活，Codex 审查")
    add_logo(sl)
    add_footer(sl, 8, TOTAL_SLIDES)
    add_code_block(sl, make_multi_agent(), top=1008000,
                   height=4500000, font_size=Pt(10))
    add_highlight_box(sl,
        "为什么用两个？一是交叉验证互相卷；"
        "二是省钱——Claude 太贵了",
        top=5700000, height=450000)

    # SLIDE 9: Perf story
    sl = add_blank_slide(prs)
    add_title(sl, "性能优化实战：寄存器固定映射的故事")
    add_logo(sl)
    add_footer(sl, 9, TOTAL_SLIDES)
    add_code_block(sl, make_perf_story(), top=1008000,
                   height=4200000, font_size=Pt(11))
    add_highlight_box(sl,
        "启示：优化方向要人类判断，Agent 负责工程落地和实验验证",
        top=5400000, height=450000)

    # SLIDE 10: Perf analysis
    sl = add_blank_slide(prs)
    add_title(sl, "Agent 怎么做性能分析")
    add_logo(sl)
    add_footer(sl, 10, TOTAL_SLIDES)
    add_code_block(sl, make_perf_analysis(), top=1008000,
                   height=3000000, font_size=Pt(11))
    add_code_block(sl, [
        "Code Quality Checks",
        ". trace log: IR translation status",
        ". host insn count: 1 Guest -> N Host?",
        ". register spill frequency",
    ], top=4200000, height=900000, font_size=Pt(12))
    add_quote(sl,
        "「一整套流程几乎不用人参与。它自己就能找出来。"
        "这体现出通用智力强不强——"
        "它的语料库里能嗅到什么是好代码。」",
        top=5300000, height=600000)

    # SLIDE 11: Benchmark
    sl = add_blank_slide(prs)
    add_title(sl, "性能实测与诚实声明")
    add_logo(sl)
    add_footer(sl, 11, TOTAL_SLIDES)
    add_code_block(sl, [
        "dhrystone 200K iterations",
        "(RISC-V -> x86-64, linux-user)",
        "",
        "QEMU   ########################################  0.4~0.5s",
        "tcg-rs ################################          0.32s",
        "                                        ^",
        "                             observed 20%~30% faster",
    ], top=1008000, height=1600000, font_size=Pt(13))
    add_table(sl, ["优化类别", "估算贡献", "关键技术"], [
        ["执行循环", "~8-10%",
         "next_tb_hint + exit_target cache"],
        ["Guest 内存访问", "~8-10%",
         "No TLB, direct [R14+addr]"],
        ["数据结构 + 其他", "~8-13%",
         "Vec storage + zero-cost abstraction"],
    ], top=2800000, row_height=380000)
    add_code_block(sl, [
        "Honest Disclosure",
        ". Only ~30% features of original implemented",
        ". Less features -> shorter path -> expected",
        ". Data value: validates optimization direction",
    ], top=4600000, height=900000, font_size=Pt(13))

    # SLIDE 12: Difftest
    sl = add_blank_slide(prs)
    add_title(sl, "差分测试：不信 AI，信参考实现")
    add_logo(sl)
    add_footer(sl, 12, TOTAL_SLIDES)
    add_code_block(sl, make_difftest(), top=1008000,
                   height=4200000, font_size=Pt(11))
    add_highlight_box(sl,
        "核心原则：QEMU 20 年打磨 >> 5 天 AI 生成 "
        "→ 以 QEMU 为 ground truth",
        top=5400000, height=450000)

    # SLIDE 13: AI boundary (left-right layout)
    sl = add_blank_slide(prs)
    add_title(sl, "AI 的边界与人类核心价值")
    add_logo(sl)
    add_footer(sl, 13, TOTAL_SLIDES)

    # Left column: Peak at First Sight + Human core value
    LEFT = 540000
    LEFT_W = 5100000
    add_code_block(sl, [
        '"Peak at First Sight"',
        "",
        "AI code gives great first",
        "impression -- clean style,",
        "clear structure",
        "",
        "Looks like mid-level or",
        "near-senior engineer",
        "",
        "But deeper review reveals",
        "internal issues",
        "-> First look IS the",
        "   best look",
    ], left=LEFT, top=1008000,
       width=LEFT_W, height=2500000, font_size=Pt(12))
    add_code_block(sl, [
        "Human core value =",
        "Define problems +",
        "Verify results",
        "",
        ". Knowledge >= AI to",
        "  collaborate effectively",
        '. AI forces learning --',
        '  or you get "fooled"',
        ". Production quality needs",
        "  human depth; AI is a",
        "  tool, not a replacement",
    ], left=LEFT, top=3700000,
       width=LEFT_W, height=2400000, font_size=Pt(12))

    # Right column: table
    RIGHT = 5900000
    RIGHT_W = 5400000
    add_table(sl, ["AI 擅长", "AI 不擅长"], [
        ["C → Rust 对照翻译", "性能优化方向判断"],
        ["批量生成 184 条", "上下文变长后"],
        ["指令函数", "遗忘细节"],
        ["测试编写 +", "多线程竞态"],
        ["边界覆盖", "条件调试"],
        ["设计文档 +", "全局架构"],
        ["代码分析", "最优解"],
    ], left=RIGHT, top=1008000,
       width=RIGHT_W, row_height=380000)

    # SLIDE 14: Expert x Model (left-right layout)
    sl = add_blank_slide(prs)
    add_title(sl, "专家经验 × Agent 智力")
    add_logo(sl)
    add_footer(sl, 14, TOTAL_SLIDES)

    # Left column: lever analogy
    add_code_block(sl, make_expert_model(),
                   left=540000, top=1008000,
                   width=5100000, height=4800000,
                   font_size=Pt(10))

    # Right column: dynamic roles table
    add_code_block(sl, make_dynamic_roles(),
                   left=5900000, top=1008000,
                   width=5400000, height=4800000,
                   font_size=Pt(9))

    add_highlight_box(sl,
        "模型越强，你不需要知道更多，但你需要判断得更准",
        top=6050000, height=400000)

    # SLIDE 15: Three practices
    sl = add_blank_slide(prs)
    add_title(sl, "三个关键实践")
    add_logo(sl)
    add_footer(sl, 15, TOTAL_SLIDES)
    add_code_block(sl, make_three_practices(), top=1008000,
                   height=4200000, font_size=Pt(10))
    add_quote(sl,
        "这三个实践不是 tcg-rs 专属的，"
        "任何系统级 AI 辅助开发都可以借鉴",
        top=5400000, height=500000)

    # SLIDE 16: Humanize
    sl = add_blank_slide(prs)
    add_title(sl, "从手动到自动：Humanize 插件")
    add_logo(sl)
    add_footer(sl, 16, TOTAL_SLIDES)
    add_code_block(sl, make_humanize(), top=1008000,
                   height=3600000, font_size=Pt(10))
    add_highlight_box(sl,
        "核心理念：迭代优于完美 · 独立审查防盲点 · "
        "多层防护防漂移",
        top=4800000, height=450000)
    add_quote(sl,
        "开源地址：github.com/humania-org/humanize",
        top=5500000, height=400000)

    # SLIDE 17: Data overview
    sl = add_blank_slide(prs)
    add_title(sl, "项目数据总览")
    add_logo(sl)
    add_footer(sl, 17, TOTAL_SLIDES)
    add_table(sl, ["指标", "数值", "说明"], [
        ["开发时间", "5 天", "全程 AI Agent 驱动"],
        ["总代码量", "~4.5 万行", "Rust（含生成代码和测试）"],
        ["测试数量", "816 个", "含差分测试 35 个，测试占比 50%"],
        ["Crate 数", "10 个", "分层解耦设计"],
        ["IR Opcodes", "158 个", "vs QEMU ~250（-40%）"],
        ["指令翻译", "184 条", "RV64IMAFDC 完整"],
        ["性能", "快 20%~30%", "vs QEMU（linux-user, dhrystone）"],
        ["花费", "~200 美元", "Claude + Codex"],
    ], top=1008000, row_height=400000)
    add_highlight_box(sl,
        "核心观点：AI 已具备系统级工程能力。"
        "关键不在于 AI 写代码有多强，"
        "而在于人类定义问题和验证结果的能力有多强。",
        top=4900000, height=600000)

    # SLIDE 18: Q&A
    sl = add_blank_slide(prs)
    sl.shapes.add_picture("/tmp/pptx_bg.png", 0, 0,
                          SLIDE_W, SLIDE_H)
    sl.shapes.add_picture("/tmp/pptx_logo.png",
                          Emu(720000), Emu(540000),
                          Emu(1800000), Emu(750802))
    t = tb(sl, Emu(720000), Emu(2160000),
           Emu(7200000), Emu(900000))
    set_run(t.text_frame.paragraphs[0], "谢谢！",
            size=Pt(52), bold=True, color=ACCENT)
    t = tb(sl, Emu(720000), Emu(3240000),
           Emu(7200000), Emu(540000))
    set_run(t.text_frame.paragraphs[0], "Q & A",
            size=Pt(28), bold=True, color=DARK)
    t = tb(sl, Emu(720000), Emu(4140000),
           Emu(7200000), Emu(432000))
    set_run(t.text_frame.paragraphs[0],
            "github.com/patchfx/tcg-rs",
            size=Pt(16), color=GRAY)

    prs.save("os2atc-2026.pptx")
    print(f"Done! {len(prs.slides)} slides saved.")


build()
